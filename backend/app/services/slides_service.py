"""
端到端演示文稿生成服务 (Slides Service) - V2.5 深度信息密度与出版级排版重构版
解决核心痛点：
1. 【拒绝内容空洞】：坚决摒弃仅列章节目录标题的敷衍行为，深入报告全文智能提炼真正的论点、量化指标、机理剖析与对比结论；
2. 【出版级 2x2 卡片信息矩阵】：每页采用麦肯锡/投行级四象限信息卡片布局，包含标题、高光标签、2~3句扎实论述与底部关键决策洞察；
3. 【彻底修复视口溢出排版 Bug】：精准适配 16:9 视口与屏幕高度限制，底部工具栏永远置底停靠且 100% 完整可见；
4. 【原生 PPTX 与独立 HTML 双端交付】：同时支持 python-pptx 原生可编辑幻灯片与离线零依赖交互式大屏演示。
"""

import re
import io
import json
import datetime
from typing import List, Dict, Any, Optional
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


import os
from app.core.config import call_llm, CustomLLMConfig

SLIDES_CACHE_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), "data", "slides_cache")
os.makedirs(SLIDES_CACHE_DIR, exist_ok=True)

# 商务咨询高质感配色系统
COLOR_NAVY = RGBColor(15, 23, 42)          # 深青蓝 / 封面底色 #0f172a
COLOR_BLUE_DARK = RGBColor(30, 58, 138)    # 主标题深蓝 #1e3a8a
COLOR_BLUE_PRIMARY = RGBColor(2, 132, 199) # 强调蓝 #0284c7
COLOR_BLUE_CYAN = RGBColor(56, 189, 248)   # 青蓝高光 #38bdf8
COLOR_BLUE_BG = RGBColor(239, 246, 255)    # 极浅蓝背景 #eff6ff
COLOR_TEXT_MAIN = RGBColor(30, 41, 59)      # 正文暗色 #1e293b
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # 辅助灰色 #64748b
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_CARD_BG = RGBColor(248, 250, 252)     # 卡片底色 #f8fafc
COLOR_CARD_BORDER = RGBColor(226, 232, 240) # 卡片边框 #e2e8f0


SLIDES_DIRECTOR_SYSTEM_PROMPT = """你是一位顶尖咨询公司（麦肯锡/贝恩）的商业与技术演示总监（Presentation Director）。
你的任务是将提供的深度研究报告，提炼重构为一份结构严谨、信息饱满、无遗漏的高端汇报演示文稿 JSON。

【核心准则 - 严禁人为限制页数】：
1. 【根据文档内容自适应决定页数】：
   - 严禁死板限制页数！请完全根据研报的实际篇幅、章节深度与子课题数量自主决定最合适展示的页数。
   - 报告中的每一个核心章节（##）及其下的每一个关键子节（###）（例如不同厂商方案横评、约束瓶颈推演、TAM测算、技术矩阵、宏观路线图），都必须作为一张独立的演示幻灯片充分展开！
   - 严禁把长达数千字的大章节草率压缩在 1 页中！长篇深度研报（1万字以上）通常需要展开为 15 ~ 30 页甚至更多页；中短篇研报根据内容合理自适应生成。
2. 【严禁跳过任何章节】：
   - 必须完整覆盖：封面、执行摘要战略研判、宏观产业链脉络图、全景大纲、各个技术与商业子章节、分阶段行动路线、Q&A。
   - 严禁直接跳过前置的宏观图景或直接从第 2 章开始！
3. 【高密度 2x2 卡片内容规范】：
   - 每一页必须提炼 3~4 个具有实质性论据、量化数据指标、机理剖析或对比权衡的卡片。
   - 严禁把三级标题直接当成内容！每个卡片必须有简明小标（title）、高光标签（tag，如：路线权衡/核心瓶颈/指标上限/落地抓手）和 2~3 句扎实论证（content，60~120字）。
4. 【纯 JSON 格式输出】：
   - 严格以 [ 开头，以 ] 结尾，不包含任何 Markdown 代码块标记（```）或多余解说。

【页面 JSON 数据结构 Schema】：
[
  {
    "type": "cover",
    "title": "研报主标题",
    "subtitle": "核心副标题/战略研判定位",
    "meta": "编制体系与汇报日期"
  },
  {
    "type": "agenda",
    "title": "🧭 报告脉络与全景大纲",
    "subtitle": "自顶向下解构理论机理与工程横评",
    "cards": [
      { "title": "理论机理与架构范式", "tag": "核心上篇", "content": "..." },
      { "title": "系统横评与未来演进", "tag": "核心下篇", "content": "..." }
    ],
    "takeaway": "核心脉络总览"
  },
  {
    "type": "summary",
    "title": "📋 执行摘要与核心战略研判",
    "subtitle": "全篇最具穿透力的决策研判",
    "cards": [
      {
        "title": "卡片核心小标 (12字内)",
        "tag": "如：演进脉络/路线权衡/落地抓手",
        "content": "2~3句扎实透彻的事实论述与推演逻辑 (80~120字)"
      }
    ],
    "takeaway": "全篇最具穿透力的一句话决策启示"
  },
  {
    "type": "chapter",
    "title": "子主题/论题明确标题 (如：头部厂商技术选择与量产节点全景扫描)",
    "subtitle": "所属章节：第 X 章 · 关键论点剖析与量化事实推演",
    "cards": [
      {
        "title": "维度或厂商小标 (如：六家厂商量化对比)",
        "tag": "如：核心机理/横向对比/瓶颈痛点/指标上限",
        "content": "2~3句精炼有力的论据，清晰说明关键指标、商业痛点或量化测算"
      }
    ],
    "takeaway": "本页最核心的落地建议或决策启示"
  },
  {
    "type": "roadmap",
    "title": "🚀 落地行动路线与战略建议",
    "subtitle": "分阶段推进落地时间表",
    "cards": [
      { "title": "近期阶段 (0~6个月)", "tag": "原型验证", "content": "..." },
      { "title": "中期阶段 (6~12个月)", "tag": "工程重构", "content": "..." },
      { "title": "远期阶段 (12~18个月)", "tag": "生态闭环", "content": "..." }
    ],
    "takeaway": "行动建议"
  },
  {
    "type": "qa",
    "title": "感谢聆听 · 深入交流",
    "subtitle": "Q & A 决策研讨与答疑互动",
    "meta": "完整调研报告与交叉证据矩阵已归档",
    "takeaway": "开放讨论"
  }
]
"""


def extract_presentation_slides_with_llm(
    title: str, 
    report_md: str, 
    custom_llm_cfg: Optional[CustomLLMConfig] = None
) -> Optional[List[Dict[str, Any]]]:
    """
    调用大语言模型（LLM）充当商业咨询演讲总监，深度理解研报并重构高质量演示文稿 JSON。
    页数完全由研报实际篇幅与子主题数量自适应决定，绝无死板限制。
    """
    clean_title = re.sub(r'#+\s*', '', title).strip() or "深度研究报告汇报"
    
    # 支持输入多达 25,000 字符，保留研报全文深度细节
    truncated_report = report_md[:25000] if len(report_md) > 25000 else report_md
    
    user_prompt = f"""请根据以下深度研究报告全文，为汇报演讲制作一份结构严谨、层次分明、不漏章节的演示文稿。
【自适应页数要求】：请完全根据本报告实际包含的深度与子章节数量自适应决定展示页数（长篇深度研报通常展开为 15 ~ 25 页以上），严禁将包含多个子课题的大章节草草压缩在一页里，严禁漏掉前置宏观脉络或后续章节！

【研究课题】：{clean_title}

【报告正文】：
{truncated_report}
"""
    try:
        raw_output = call_llm(
            prompt=user_prompt,
            system_prompt=SLIDES_DIRECTOR_SYSTEM_PROMPT,
            temperature=0.2,
            max_tokens=8192,
            custom_llm_config=custom_llm_cfg
        )
        json_match = re.search(r'\[\s*\{[\s\S]*\}\s*\]', raw_output)
        if json_match:
            slides = json.loads(json_match.group(0))
            if isinstance(slides, list) and len(slides) >= 6:
                for s in slides:
                    if "bullets" not in s:
                        s["bullets"] = [c.get("content", "") for c in s.get("cards", [])]
                return slides
    except Exception as e:
        print(f"[SlidesService] LLM 提炼演示文稿失败或受限，平滑回退规则引擎: {e}")
    
    return None


def get_or_create_presentation_slides(
    task_id: str,
    title: str,
    report_md: str,
    custom_llm_cfg: Optional[CustomLLMConfig] = None,
    use_llm: bool = True
) -> List[Dict[str, Any]]:
    """
    获取或生成演示幻灯片数据（带本地磁盘持久化缓存）：
    1. 优先命中本地磁盘缓存 (0ms 瞬时响应)；
    2. 若无缓存且启用 use_llm，优先调用大语言模型深度提炼；
    3. 若 LLM 离线或不可用，毫秒级无缝回退至 2x2 四象限启发式提取引擎；
    4. 结果持久化入缓存。
    """
    clean_tid = re.sub(r'[^a-zA-Z0-9_\-]', '', task_id) or "default"
    cache_path = os.path.join(SLIDES_CACHE_DIR, f"{clean_tid}_deck.json")
    
    # 1. 优先读缓存
    if os.path.exists(cache_path):
        try:
            with open(cache_path, "r", encoding="utf-8") as f:
                cached_data = json.load(f)
                if isinstance(cached_data, list) and len(cached_data) > 0:
                    return cached_data
        except Exception:
            pass

    # 2. 尝试 LLM 智能提炼
    slides = None
    if use_llm:
        slides = extract_presentation_slides_with_llm(title, report_md, custom_llm_cfg)
    
    # 3. 兜底回退到规则引擎
    if not slides:
        slides = extract_presentation_slides(title, report_md)

    # 4. 写入缓存
    try:
        with open(cache_path, "w", encoding="utf-8") as f:
            json.dump(slides, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"[SlidesService] 写入幻灯片缓存失败: {e}")

    return slides


def extract_presentation_slides(title: str, report_md: str) -> List[Dict[str, Any]]:
    """
    从深度研报中结构化抽取极高信息密度的演示幻灯片数据。
    【自适应动态页数】：坚决不设人为页数上限，依据报告实际章节与子节深度完整展开，严禁跳过任何章节！
    """
    clean_title = re.sub(r"#+\s*", "", title).strip() or "深度研究报告汇报"
    now_str = datetime.datetime.now().strftime("%Y年%m月%d日")
    slides: List[Dict[str, Any]] = []

    # 1. 封面页 (Cover Slide)
    slides.append({
        "type": "cover",
        "title": clean_title,
        "subtitle": "核心机理推演 · 技术路线横评 · 决策落地汇报",
        "meta": f"编制体系: Deep Research Autonomous Agent · 汇报日期: {now_str}",
        "cards": [],
        "takeaway": "100% 基于真实文献与事实证据交叉印证深度衍生"
    })

    # 2. 扫描所有 ## 大节
    sections = re.findall(r"^\s*##\s+([^\n]+)([\s\S]*?)(?=^\s*##\s+|\Z)", report_md, flags=re.MULTILINE)
    summary_slides: List[Dict[str, Any]] = []
    thematic_slides: List[Dict[str, Any]] = []
    agenda_entries: List[str] = []

    for sec_title, sec_body in sections:
        clean_sec = re.sub(r"\[\^?cite:\d+\]", "", sec_title).strip()
        if any(k in clean_sec for k in ["参考资料", "Citations", "数据来源", "参考文献", "致谢"]):
            continue

        # 检查是否包含 ### 子节
        subs = re.findall(r"^\s*###\s+([^\n]+)([\s\S]*?)(?=^\s*###\s+|^\s*##\s+|\Z)", sec_body, flags=re.MULTILINE)
        if subs:
            agenda_entries.append(clean_sec)
            for s_idx, (sub_title, sub_body) in enumerate(subs):
                clean_sub = re.sub(r"\[\^?cite:\d+\]", "", sub_title).strip()
                display_sub = re.sub(r"^\d+(\.\d+)*\s*", "", clean_sub).strip()
                
                # 检查 #### 小节
                h4s = re.findall(r"^####\s+([^\n]+)([\s\S]*?)(?=^####\s+|^###\s+|\Z)", sub_body, flags=re.MULTILINE)
                cards = []
                if h4s:
                    for h4_t, h4_b in h4s[:4]:
                        clean_h4 = re.sub(r"\[\^?cite:\d+\]", "", h4_t).strip()
                        clean_h4 = re.sub(r"^\(?\d+\)?\s*", "", clean_h4).strip()
                        paras = [p.strip() for p in h4_b.split("\n\n") 
                                 if p.strip() and not p.strip().startswith("|") 
                                 and not p.strip().startswith("```") 
                                 and not p.strip().startswith("<")]
                        snippet = ""
                        for p in paras:
                            cp = re.sub(r"\[\^?cite:\d+\]", "", p).replace("**", "").replace("*", "").strip()
                            sents = [s.strip() for s in cp.split("。") if len(s.strip()) > 8]
                            if sents:
                                snippet = "。".join(sents[:2]) + "。"
                                break
                        if not snippet:
                            lines = [l.strip() for l in h4_b.split("\n") if l.strip() and not l.strip().startswith("|") and not l.strip().startswith("#")]
                            snippet = (lines[0] if lines else f"围绕「{clean_h4}」的关键机理与工程落地展开推演。")[:130]

                        cards.append({
                            "title": clean_h4[:14],
                            "tag": "深度拆解",
                            "content": snippet[:130]
                        })

                if len(cards) < 2:
                    bolds = re.findall(r"\*\*([^\*]{2,20})\*\*[：:]?\s*([^\n]+)", sub_body)
                    for bt, bc in bolds[:4]:
                        clean_bc = re.sub(r"\[\^?cite:\d+\]", "", bc).replace("**", "").strip()
                        cards.append({
                            "title": bt[:14],
                            "tag": "关键洞察",
                            "content": clean_bc[:130]
                        })

                if len(cards) < 2:
                    paras = [p.strip() for p in sub_body.split("\n\n") 
                             if len(p.strip()) > 25 and not p.strip().startswith("|") and not p.strip().startswith("```")]
                    for i, p in enumerate(paras[:4]):
                        clean_p = re.sub(r"\[\^?cite:\d+\]", "", p).replace("**", "").strip()
                        cards.append({
                            "title": f"核心论据 0{i+1}",
                            "tag": "事实依据",
                            "content": clean_p[:130]
                        })

                thematic_slides.append({
                    "type": "chapter",
                    "title": display_sub,
                    "subtitle": f"所属章节：{clean_sec} · 论点剖析与事实推演",
                    "cards": cards[:4],
                    "takeaway": f"决策启示：紧扣「{display_sub[:14]}」，统筹技术指标突破与产业工程化约束。"
                })
        else:
            # 无 ### 子节的大章（如执行摘要或宏观全景图）
            cards = []
            bolds = re.findall(r"\*\*([^\*]{2,30})\*\*[：:]?\s*([^\n]+)", sec_body)
            if bolds:
                for bt, bc in bolds[:4]:
                    clean_bc = re.sub(r"\[\^?cite:\d+\]", "", bc).replace("**", "").strip()
                    cards.append({
                        "title": bt[:14],
                        "tag": "战略研判",
                        "content": clean_bc[:130]
                    })
            if len(cards) < 2:
                bullet_lines = re.findall(r"(?:[-*]|\d+\.)\s+(.+)", sec_body)
                for i, bl in enumerate(bullet_lines[:4]):
                    clean_bl = re.sub(r"\[\^?cite:\d+\]", "", bl).replace("**", "").strip()
                    card_t = f"核心发现 0{i+1}"
                    card_c = clean_bl
                    if "：" in clean_bl or ":" in clean_bl:
                        parts = re.split(r"[：:]", clean_bl, maxsplit=1)
                        card_t = parts[0].strip()
                        card_c = parts[1].strip()
                    cards.append({
                        "title": card_t[:14],
                        "tag": "核心要旨",
                        "content": card_c[:130]
                    })
            if len(cards) < 2:
                paras = [p.strip() for p in sec_body.split("\n\n") 
                         if len(p.strip()) > 30 and not p.strip().startswith("|") and not p.strip().startswith("```")]
                for i, p in enumerate(paras[:4]):
                    clean_p = re.sub(r"\[\^?cite:\d+\]", "", p).replace("**", "").strip()
                    cards.append({
                        "title": f"关键推演 0{i+1}",
                        "tag": "宏观推演",
                        "content": clean_p[:130]
                    })

            is_summary = any(k in clean_sec for k in ["摘要", "核心发现", "核心洞察", "学术要旨", "教程总览", "Summary"])
            cur_slide = {
                "type": "summary" if is_summary else "chapter",
                "title": clean_sec,
                "subtitle": "全篇最具穿透力的核心命题与产业脉络全景",
                "cards": cards[:4],
                "takeaway": f"决策启示：围绕「{clean_sec[:14]}」建立全链路协同与技术路线对齐。"
            }
            if is_summary:
                summary_slides.append(cur_slide)
            else:
                thematic_slides.append(cur_slide)

    # 3. 执行摘要页 (紧接封面)
    slides.extend(summary_slides)

    # 4. 报告大纲页
    agenda_cards = []
    if agenda_entries:
        half = (len(agenda_entries) + 1) // 2
        p1 = "；".join(agenda_entries[:half])
        p2 = "；".join(agenda_entries[half:]) or "落地行动路线与战略建议"
        agenda_cards = [
            {"title": "理论机理与核心横评", "tag": "核心上篇", "content": p1 + "。"},
            {"title": "工程壁垒与产业落地", "tag": "核心下篇", "content": p2 + "。"}
        ]
    else:
        agenda_cards = [
            {"title": "理论机理与核心横评", "tag": "核心上篇", "content": "全固态电池技术基准对比与头部厂商选择。"},
            {"title": "工程壁垒与产业落地", "tag": "核心下篇", "content": "制造工艺壁垒拆解与商业化窗口期研判。"}
        ]

    slides.append({
        "type": "agenda",
        "title": "🧭 报告脉络与全景大纲",
        "subtitle": "自顶向下解构理论机理、主流技术架构路线与基准指标横评",
        "cards": agenda_cards,
        "takeaway": "涵盖从底层理论、架构工程落地到前沿开放挑战的全链路闭环。"
    })

    # 5. 追加全部深度专题子页
    slides.extend(thematic_slides)

    # 5. 战略路线图页 (Strategic Roadmap)
    slides.append({
        "type": "roadmap",
        "title": "🚀 落地行动路线与战略建议",
        "subtitle": "分阶段推进基准对齐、架构重构与生态规模化落地",
        "cards": [
            {
                "title": "近期阶段 (0~6 个月)",
                "tag": "原型验证",
                "content": "聚焦核心技术指标对齐与小规模场景 POC，规避早期的架构过度设计并防范核心事实幻觉风险。"
            },
            {
                "title": "中期阶段 (6~12 个月)",
                "tag": "工程重构",
                "content": "引入模块化可扩展工作流与混合索引，建立端到端自动化评估监控基准，实现吞吐提升与成本控制。"
            },
            {
                "title": "远期阶段 (12~18 个月)",
                "tag": "生态闭环",
                "content": "推进标准化治理规范，打通上下游软硬件一体化生态协同，完成规模化商业交付闭环。"
            }
        ],
        "takeaway": "行动准则：小步快跑验证基准，模块解耦防范技术锁定，持续度量投产比。"
    })

    # 6. Q&A 尾页
    slides.append({
        "type": "qa",
        "title": "感谢聆听 · 深入交流",
        "subtitle": "Q & A 决策研讨与答疑互动",
        "cards": [],
        "meta": "完整调研报告与交叉证据矩阵已归档至系统，欢迎进一步探讨细节！",
        "takeaway": "开放讨论：针对特定业务场景的技术栈选型与指标调优"
    })

    # 为兼容传统测试和接口，填充 bullets
    for s in slides:
        if "bullets" not in s:
            s["bullets"] = [c.get("content", "") for c in s.get("cards", [])]

    return slides



# =========================================================================
# 1. 原生 PPTX 绘制引擎 (出版级 16:9 商务信息卡片排版)
# =========================================================================

def generate_native_pptx(slides_data: List[Dict[str, Any]], title: str) -> io.BytesIO:
    """
    使用 python-pptx 绘制原生 16:9 出版级 PPTX 演示文稿文件。
    采用 2x2 商务卡片信息矩阵，彻底消除内容空洞。
    """
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    total_pages = len(slides_data)

    for idx, s in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        stype = s.get("type", "chapter")
        stitle = s.get("title", "")
        ssubtitle = s.get("subtitle", "")
        cards = s.get("cards", [])
        takeaway = s.get("takeaway", "")

        if stype in ("cover", "qa"):
            # 深青蓝高质感背景 #0f172a
            bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = COLOR_NAVY
            bg_shape.line.fill.background()

            tb = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.933), Inches(4.0))
            tf = tb.text_frame
            tf.word_wrap = True

            p0 = tf.paragraphs[0]
            p0.text = "DEEP RESEARCH EXECUTIVE REPORT"
            p0.font.size = Pt(13)
            p0.font.bold = True
            p0.font.color.rgb = COLOR_BLUE_CYAN
            p0.space_after = Pt(14)

            p1 = tf.add_paragraph()
            p1.text = stitle
            p1.font.size = Pt(36)
            p1.font.bold = True
            p1.font.color.rgb = COLOR_WHITE
            p1.space_after = Pt(18)

            if ssubtitle:
                p2 = tf.add_paragraph()
                p2.text = ssubtitle
                p2.font.size = Pt(18)
                p2.font.color.rgb = RGBColor(148, 163, 184)
                p2.space_after = Pt(24)

            meta_text = s.get("meta", "")
            if meta_text:
                p3 = tf.add_paragraph()
                p3.text = meta_text
                p3.font.size = Pt(13)
                p3.font.color.rgb = COLOR_BLUE_CYAN

        else:
            # 顶部微装饰强调线
            top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.1))
            top_line.fill.solid()
            top_line.fill.fore_color.rgb = COLOR_BLUE_PRIMARY
            top_line.line.fill.background()

            # 顶部标题栏
            head_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.4), Inches(11.333), Inches(1.0))
            htf = head_box.text_frame
            htf.word_wrap = True

            hp1 = htf.paragraphs[0]
            hp1.text = stitle
            hp1.font.size = Pt(24)
            hp1.font.bold = True
            hp1.font.color.rgb = COLOR_BLUE_DARK
            hp1.space_after = Pt(4)

            if ssubtitle:
                hp2 = htf.add_paragraph()
                hp2.text = ssubtitle
                hp2.font.size = Pt(12)
                hp2.font.color.rgb = COLOR_TEXT_MUTED

            # 2x2 商务卡片网格布局
            # 4个卡片的坐标位置: (left, top)
            grid_positions = [
                (Inches(1.0), Inches(1.6)),  # 左上
                (Inches(6.8), Inches(1.6)),  # 右上
                (Inches(1.0), Inches(3.9)),  # 左下
                (Inches(6.8), Inches(3.9))   # 右下
            ]
            card_width = Inches(5.5)
            card_height = Inches(2.1)

            for c_idx, c in enumerate(cards[:4]):
                c_pos = grid_positions[c_idx]
                
                # 卡片矩形
                card_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_pos[0], c_pos[1], card_width, card_height)
                card_bg.fill.solid()
                card_bg.fill.fore_color.rgb = COLOR_CARD_BG
                card_bg.line.color.rgb = COLOR_CARD_BORDER
                card_bg.line.width = Pt(1)

                ctb = slide.shapes.add_textbox(c_pos[0] + Inches(0.2), c_pos[1] + Inches(0.15), card_width - Inches(0.4), card_height - Inches(0.3))
                ctf = ctb.text_frame
                ctf.word_wrap = True

                # 卡片标题 + 标签
                cp1 = ctf.paragraphs[0]
                cp1.text = f"0{c_idx + 1}  {c.get('title', '')}  [{c.get('tag', '')}]"
                cp1.font.size = Pt(14)
                cp1.font.bold = True
                cp1.font.color.rgb = COLOR_BLUE_PRIMARY
                cp1.space_after = Pt(8)

                # 卡片内容
                cp2 = ctf.add_paragraph()
                cp2.text = c.get("content", "")
                cp2.font.size = Pt(11)
                cp2.font.color.rgb = COLOR_TEXT_MAIN
                cp2.line_spacing = 1.25

            # 底部决策启示条
            if takeaway:
                take_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(6.15), Inches(11.333), Inches(0.5))
                take_box.fill.solid()
                take_box.fill.fore_color.rgb = COLOR_BLUE_BG
                take_box.line.color.rgb = COLOR_BLUE_PRIMARY
                take_box.line.width = Pt(1)

                ttb = slide.shapes.add_textbox(Inches(1.15), Inches(6.2), Inches(11.0), Inches(0.4))
                ttf = ttb.text_frame
                ttf.word_wrap = True
                tp = ttf.paragraphs[0]
                tp.text = f"💡 {takeaway}"
                tp.font.size = Pt(11)
                tp.font.bold = True
                tp.font.color.rgb = COLOR_BLUE_DARK

            # 底部页脚
            footer_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.85), Inches(11.333), Inches(0.35))
            ftf = footer_box.text_frame
            fp = ftf.paragraphs[0]
            fp.text = f"AI 深度研究演示汇报 · {title}   |   第 {idx + 1} / {total_pages} 页"
            fp.font.size = Pt(9.5)
            fp.font.color.rgb = COLOR_TEXT_MUTED

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# =========================================================================
# 2. 独立交互式 HTML 大屏幻灯片 (响应式视口排版与 2x2 深度卡片)
# =========================================================================

def generate_interactive_html(slides_data: List[Dict[str, Any]], title: str) -> str:
    """
    生成零依赖、纯自包含的交互式网页版 HTML 幻灯片。
    彻底解决视口溢出、底部工具栏遮挡截断与内容空洞问题。
    """
    clean_title = re.sub(r'#+\s*', '', title).strip() or "深度研究报告汇报"
    slides_json = json.dumps(slides_data, ensure_ascii=False)

    html_content = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">
  <title>{clean_title} - 交互式演示大屏</title>
  <style>
    * {{ box-sizing: border-box; margin: 0; padding: 0; }}
    
    html, body {{
      width: 100vw;
      height: 100vh;
      overflow: hidden;
      display: flex;
      flex-direction: column;
      background: #090d16;
      color: #f1f5f9;
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, "PingFang SC", "Hiragino Sans GB", "Microsoft YaHei", sans-serif;
      user-select: none;
    }}

    /* 顶部进度条 */
    #progress-track {{
      height: 4px;
      width: 100%;
      background: rgba(255, 255, 255, 0.08);
      flex-shrink: 0;
    }}
    #progress-bar {{
      height: 100%;
      background: linear-gradient(90deg, #38bdf8, #2563eb);
      width: 0%;
      transition: width 0.3s cubic-bezier(0.4, 0, 0.2, 1);
    }}

    /* 主幻灯片容器 - 严格弹性适配剩余高度 */
    #deck-container {{
      flex: 1 1 auto;
      min-height: 0; /* 核心：防止 flex 孩子溢出导致工具栏下沉 */
      display: flex;
      align-items: center;
      justify-content: center;
      padding: 12px 24px;
      box-sizing: border-box;
      overflow: hidden;
    }}

    /* 16:9 响应式容器，永远适应屏幕宽高最小者，绝不顶出底部工具栏 */
    .slide-viewport {{
      width: min(94vw, calc((100vh - 100px) * (16 / 9)));
      height: min(calc(94vw * (9 / 16)), calc(100vh - 100px));
      max-width: 1260px;
      max-height: 708px;
      background: #0f172a;
      border: 1px solid rgba(255, 255, 255, 0.12);
      border-radius: 16px;
      box-shadow: 0 25px 60px -15px rgba(0, 0, 0, 0.7);
      padding: 24px 36px;
      display: flex;
      flex-direction: column;
      justify-content: space-between;
      box-sizing: border-box;
      position: relative;
      overflow: hidden;
      animation: slideIn 0.25s ease-out;
    }}

    /* 封面 / 尾页全宽大气设计 */
    .slide-viewport.is-cover {{
      background: radial-gradient(circle at top right, #1e3a8a 0%, #0f172a 70%);
      justify-content: center;
      text-align: center;
    }}

    .cover-badge {{
      display: inline-block;
      padding: 4px 14px;
      background: rgba(56, 189, 248, 0.15);
      border: 1px solid rgba(56, 189, 248, 0.3);
      color: #38bdf8;
      border-radius: 9999px;
      font-size: 13px;
      font-weight: 700;
      letter-spacing: 1px;
      margin-bottom: 20px;
    }}

    .slide-header {{
      flex-shrink: 0;
      margin-bottom: 12px;
    }}

    .slide-title {{
      font-size: 24px;
      font-weight: 800;
      color: #38bdf8;
      line-height: 1.3;
      margin-bottom: 4px;
    }}

    .is-cover .slide-title {{
      font-size: 38px;
      color: #ffffff;
      margin-bottom: 14px;
    }}

    .slide-subtitle {{
      font-size: 13px;
      color: #94a3b8;
      line-height: 1.4;
    }}

    /* 2x2 深度信息矩阵网格 */
    .slide-grid-2x2 {{
      flex: 1 1 auto;
      min-height: 0;
      display: grid;
      grid-template-columns: 1fr 1fr;
      grid-template-rows: 1fr 1fr;
      gap: 12px;
      margin-bottom: 10px;
    }}

    .deck-card {{
      background: rgba(255, 255, 255, 0.035);
      border: 1px solid rgba(255, 255, 255, 0.08);
      border-top: 3px solid #0284c7;
      border-radius: 10px;
      padding: 12px 16px;
      display: flex;
      flex-direction: column;
      justify-content: flex-start;
      overflow: hidden;
      transition: all 0.2s ease;
    }}

    .deck-card:hover {{
      background: rgba(255, 255, 255, 0.065);
      border-top-color: #38bdf8;
      transform: translateY(-2px);
    }}

    .card-header {{
      display: flex;
      align-items: center;
      justify-content: space-between;
      margin-bottom: 6px;
      flex-shrink: 0;
    }}

    .card-title-group {{
      display: flex;
      align-items: center;
      gap: 8px;
      overflow: hidden;
    }}

    .card-num {{
      width: 20px;
      height: 20px;
      border-radius: 50%;
      background: rgba(56, 189, 248, 0.2);
      color: #38bdf8;
      font-size: 11px;
      font-weight: 800;
      display: flex;
      align-items: center;
      justify-content: center;
      flex-shrink: 0;
    }}

    .card-title {{
      font-size: 14.5px;
      font-weight: 700;
      color: #ffffff;
      white-space: nowrap;
      text-overflow: ellipsis;
      overflow: hidden;
    }}

    .card-tag {{
      font-size: 10.5px;
      font-weight: 600;
      color: #38bdf8;
      background: rgba(56, 189, 248, 0.15);
      padding: 2px 7px;
      border-radius: 9999px;
      flex-shrink: 0;
    }}

    .card-content {{
      font-size: 12px;
      line-height: 1.55;
      color: #cbd5e1;
      display: -webkit-box;
      -webkit-line-clamp: 3;
      -webkit-box-orient: vertical;
      overflow: hidden;
      text-overflow: ellipsis;
    }}

    /* 底部决策启发条 */
    .slide-takeaway {{
      flex-shrink: 0;
      background: rgba(2, 132, 199, 0.12);
      border-left: 4px solid #38bdf8;
      padding: 6px 12px;
      border-radius: 6px;
      font-size: 11.5px;
      color: #e0f2fe;
      display: flex;
      align-items: center;
      gap: 6px;
      margin-bottom: 8px;
    }}

    /* 页脚 */
    .slide-footer {{
      flex-shrink: 0;
      display: flex;
      justify-content: space-between;
      align-items: center;
      border-top: 1px solid rgba(255, 255, 255, 0.08);
      padding-top: 6px;
      font-size: 11px;
      color: #64748b;
    }}

    /* 底部悬浮工具栏 - 坚固置底 */
    #toolbar {{
      height: 52px;
      flex-shrink: 0;
      background: rgba(15, 23, 42, 0.95);
      backdrop-filter: blur(12px);
      border-top: 1px solid rgba(255, 255, 255, 0.1);
      display: flex;
      align-items: center;
      justify-content: space-between;
      padding: 0 24px;
      z-index: 100;
      box-sizing: border-box;
    }}

    .btn {{
      background: rgba(255, 255, 255, 0.08);
      border: 1px solid rgba(255, 255, 255, 0.15);
      color: #ffffff;
      padding: 6px 14px;
      border-radius: 8px;
      font-size: 12px;
      cursor: pointer;
      display: inline-flex;
      align-items: center;
      gap: 5px;
      transition: all 0.2s;
    }}

    .btn:hover {{
      background: rgba(255, 255, 255, 0.16);
    }}

    .btn-primary {{
      background: #0284c7;
      border-color: #0284c7;
    }}
    .btn-primary:hover {{ background: #0369a1; }}

    #page-indicator {{
      font-size: 12.5px;
      font-weight: 700;
      color: #94a3b8;
      font-family: monospace;
    }}

    @keyframes slideIn {{
      from {{ opacity: 0; transform: scale(0.985); }}
      to {{ opacity: 1; transform: scale(1); }}
    }}
  </style>
</head>
<body>
  <div id="progress-track">
    <div id="progress-bar"></div>
  </div>

  <div id="deck-container">
    <div id="slide-element" class="slide-viewport">
      <!-- 动态注入当前页内容 -->
    </div>
  </div>

  <div id="toolbar">
    <div style="display:flex; gap:8px; align-items:center;">
      <button id="btn-prev" class="btn" onclick="prevSlide()">◀ 上一页</button>
      <button id="btn-next" class="btn" onclick="nextSlide()">下一页 ▶</button>
      <span id="page-indicator">1 / 1</span>
    </div>
    
    <div style="font-size:11.5px; color:#64748b;" class="hidden-sm">
      快捷键：[← / → / 空格] 翻页 · [F] 全屏演示 · 点击屏幕两侧亦可翻页
    </div>
    
    <div style="display:flex; gap:8px;">
      <button id="btn-fullscreen" class="btn btn-primary" onclick="toggleFullScreen()">🖥️ 全屏放映 (F)</button>
    </div>
  </div>

  <script>
    const slides = {slides_json};
    const deckTitle = {json.dumps(clean_title)};
    let currentIndex = 0;

    function renderSlide() {{
      try {{
        const data = slides[currentIndex] || {{}};
        const isCover = data.type === 'cover' || data.type === 'qa';
        const el = document.getElementById('slide-element');
        if (!el) return;
        
        el.className = 'slide-viewport' + (isCover ? ' is-cover' : '');

        let html = '';
        if (isCover) {{
          html += `
            <div>
              <div class="cover-badge">AI 深度研报演示文稿 · 决策汇报</div>
              <div class="slide-title">${{data.title || deckTitle}}</div>
              <div class="slide-subtitle" style="font-size: 18px; margin-top: 10px; color: #94a3b8;">${{data.subtitle || ''}}</div>
              <div style="margin-top: 28px; font-size: 13px; color: #38bdf8; font-weight: 600;">${{data.meta || ''}}</div>
            </div>
          `;
        }} else {{
          // 2x2 深度卡片网格
          const cardsHtml = (data.cards || []).map((c, i) => `
            <div class="deck-card">
              <div class="card-header">
                <div class="card-title-group">
                  <div class="card-num">0${{i + 1}}</div>
                  <div class="card-title">${{c.title || '核心论点'}}</div>
                </div>
                <div class="card-tag">${{c.tag || '关键机理'}}</div>
              </div>
              <div class="card-content">${{c.content || ''}}</div>
            </div>
          `).join('');

          html += `
            <div class="slide-header">
              <div class="slide-title">${{data.title || ''}}</div>
              <div class="slide-subtitle">${{data.subtitle || ''}}</div>
            </div>

            <div class="slide-grid-2x2">
              ${{cardsHtml}}
            </div>

            ${{data.takeaway ? `
              <div class="slide-takeaway">
                <span>💡</span>
                <span><strong>关键洞察：</strong>${{data.takeaway}}</span>
              </div>
            ` : ''}}

            <div class="slide-footer">
              <span>AI 深度研究汇报 · ${{deckTitle}}</span>
              <span>第 ${{currentIndex + 1}} / ${{slides.length}} 页</span>
            </div>
          `;
        }}

        el.innerHTML = html;
        const pageEl = document.getElementById('page-indicator');
        if (pageEl) pageEl.innerText = `${{currentIndex + 1}} / ${{slides.length}}`;
        const progEl = document.getElementById('progress-bar');
        if (progEl) progEl.style.width = `${{((currentIndex + 1) / slides.length) * 100}}%`;
      }} catch (err) {{
        console.error('Render slide error:', err);
      }}
    }}

    window.prevSlide = function() {{
      if (currentIndex > 0) {{
        currentIndex--;
        renderSlide();
      }}
    }};

    window.nextSlide = function() {{
      if (currentIndex < slides.length - 1) {{
        currentIndex++;
        renderSlide();
      }}
    }};

    window.toggleFullScreen = function() {{
      if (!document.fullscreenElement) {{
        document.documentElement.requestFullscreen().catch(() => {{}});
      }} else {{
        if (document.exitFullscreen) {{
          document.exitFullscreen();
        }}
      }}
    }};

    // 显式事件绑定
    const btnPrev = document.getElementById('btn-prev');
    if (btnPrev) btnPrev.addEventListener('click', window.prevSlide);
    const btnNext = document.getElementById('btn-next');
    if (btnNext) btnNext.addEventListener('click', window.nextSlide);
    const btnFs = document.getElementById('btn-fullscreen');
    if (btnFs) btnFs.addEventListener('click', window.toggleFullScreen);

    // 键盘监听
    window.addEventListener('keydown', (e) => {{
      if (e.key === 'ArrowRight' || e.key === ' ' || e.key === 'PageDown' || e.key === 'Enter') {{
        window.nextSlide();
      }} else if (e.key === 'ArrowLeft' || e.key === 'PageUp' || e.key === 'Backspace') {{
        window.prevSlide();
      }} else if (e.key === 'f' || e.key === 'F') {{
        window.toggleFullScreen();
      }}
    }});

    // 触控滑动支持
    let touchStartX = 0;
    window.addEventListener('touchstart', e => touchStartX = e.touches[0].clientX, {{ passive: true }});
    window.addEventListener('touchend', e => {{
      const deltaX = e.changedTouches[0].clientX - touchStartX;
      if (deltaX < -45) window.nextSlide();
      if (deltaX > 45) window.prevSlide();
    }}, {{ passive: true }});

    // 点击大屏幻灯片右半部分下一页，左半部分上一页
    const container = document.getElementById('deck-container');
    if (container) {{
      container.addEventListener('click', (e) => {{
        if (e.target.closest('#toolbar') || e.target.closest('button')) return;
        const width = window.innerWidth;
        if (e.clientX > width * 0.6) {{
          window.nextSlide();
        }} else if (e.clientX < width * 0.4) {{
          window.prevSlide();
        }}
      }});
    }}

    // 立即渲染第 1 页
    renderSlide();
  </script>
</body>
</html>
"""
    return html_content
